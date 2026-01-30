import streamlit as st
import requests
import xml.etree.ElementTree as ET
import pandas as pd
from io import BytesIO
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
import xlsxwriter
from urllib.parse import quote_plus
import time
import urllib3
import datetime
import random
import folium
from streamlit_folium import st_folium
import re

# SSL 경고 비활성화
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# =========================================================
# [설정] UI 및 스타일
# =========================================================
st.set_page_config(page_title="부동산 원클릭 분석 Pro", page_icon="🏢", layout="centered")

st.markdown("""
    <style>
        @import url("https://cdn.jsdelivr.net/gh/orioncactus/pretendard/dist/web/static/pretendard.css");
        html, body, [class*="css"] { font-family: 'Pretendard', sans-serif; }
        .block-container { max-width: 1000px; padding: 3rem 2rem; }
        button[data-testid="stNumberInputStepDown"], button[data-testid="stNumberInputStepUp"] { display: none !important; }
        .stNumberInput label { display: none; }
          
        input[type="text"] { text-align: right !important; font-size: 19px !important; font-weight: 700 !important; color: #111 !important; background-color: #fdfdfd !important; border: 2px solid #e0e0e0 !important; }
        div[data-testid="stTextInput"] input[aria-label="주소 입력"] { text-align: left !important; }
        div[data-testid="stTextInput"] input[aria-label="공시지가"], div[data-testid="stTextInput"] input[aria-label="용도지역"] { text-align: center !important; color: #1a237e !important; }
        input[aria-label="매매금액"] { color: #D32F2F !important; font-size: 34px !important; font-weight: 900 !important; }
          
        .stButton > button { width: 100%; background-color: #1a237e; color: white; font-size: 18px; font-weight: 800; padding: 14px; border-radius: 8px; border: none; box-shadow: 0 4px 6px rgba(0,0,0,0.2); transition: all 0.3s; }
        .stButton > button:hover { background-color: #0d47a1; transform: translateY(-2px); }
          
        div[data-testid="column"] .stButton > button { background-color: transparent !important; border: none !important; color: black !important; font-weight: 900 !important; font-size: 18px !important; box-shadow: none !important; padding: 5px !important; text-align: center !important; }
        div[data-testid="column"] .stButton > button:hover { color: #333 !important; background-color: transparent !important; transform: scale(1.05); }

        .unit-price-box { background-color: #f0f2f5; border: 2px solid #d1d5db; padding: 12px; border-radius: 10px; margin-top: 10px; text-align: center; }
        .unit-price-value { font-size: 24px; font-weight: 900; color: #000; }
          
        .ai-summary-box { background-color: #fff; border: 2px solid #c5cae9; border-top: 5px solid #1a237e; padding: 30px; border-radius: 8px; margin-top: 20px; text-align: left; box-shadow: 0 10px 25px rgba(0,0,0,0.08); }
        .ai-title { font-size: 26px; font-weight: 900; color: #1a237e; margin-bottom: 25px; border-bottom: 3px solid #eee; padding-bottom: 15px; letter-spacing: -0.5px; }
          
        .link-btn { display: inline-block; width: 100%; padding: 12px; margin: 5px 0; text-align: center; border-radius: 6px; text-decoration: none; font-weight: 800; font-size: 16px; color: white !important; transition: 0.3s; box-shadow: 0 2px 4px rgba(0,0,0,0.1); }
        .naver-btn { background-color: #03C75A; } .eum-btn { background-color: #1a237e; }
        .naver-btn:hover, .eum-btn:hover { opacity: 0.8; transform: translateY(-1px); }
          
        .selected-tags { background-color: #e3f2fd; color: #0d47a1; padding: 8px 14px; border-radius: 20px; font-size: 15px; font-weight: 800; margin-right: 6px; display: inline-block; margin-bottom: 8px; border: 2px solid #90caf9; }
          
        div[data-testid="stTextInput"] input[aria-label="대지면적"], div[data-testid="stTextInput"] input[aria-label="연면적"], div[data-testid="stTextInput"] input[aria-label="건축면적"], div[data-testid="stTextInput"] input[aria-label="지상면적"] { font-size: 26px !important; font-weight: 900 !important; color: #000 !important; }

        [data-testid="stDataFrame"] { font-size: 17px !important; font-weight: 700 !important; }
        [data-testid="stElementToolbar"] { opacity: 1 !important; visibility: visible !important; font-weight: bold; }
        [data-testid="stElementToolbarButton"] { transform: scale(1.1); }
        div[data-testid="stDataFrame"] th { font-size: 16px !important; font-weight: 900 !important; color: #1a237e !important; background-color: #f8f9fa !important; }
    </style>
    """, unsafe_allow_html=True)

# =========================================================
# [필수 수정] API 키 및 주소 설정 (확인됨)
# =========================================================
VWORLD_KEY = "92DFF41C-AAAD-327C-AF08-5439410E69A4"
REFERER_URL = "https://port-0-infobd-app-mkz6091j1bce3145.sel3.cloudtype.app/"
USER_KEY = "Xl5W1ALUkfEhomDR8CBUoqBMRXphLTIB7CuTto0mjsg0CQQspd7oUEmAwmw724YtkjnV05tdEx6y4yQJCe3W0g=="

if 'zoning' not in st.session_state: st.session_state['zoning'] = ""
if 'generated_candidates' not in st.session_state: st.session_state['generated_candidates'] = [] 
if 'final_selected_insights' not in st.session_state: st.session_state['final_selected_insights'] = [] 
if 'price' not in st.session_state: st.session_state['price'] = 0
if 'addr' not in st.session_state: st.session_state['addr'] = "" 
if 'last_click_lat' not in st.session_state: st.session_state['last_click_lat'] = 0.0
if 'fetched_lp' not in st.session_state: st.session_state['fetched_lp'] = 0
if 'fetched_zoning' not in st.session_state: st.session_state['fetched_zoning'] = ""
if 'rent_roll_data' not in st.session_state: st.session_state['rent_roll_data'] = []
if 'rent_roll_init' not in st.session_state: st.session_state['rent_roll_init'] = False

def reset_analysis():
    st.session_state['generated_candidates'] = []
    st.session_state['final_selected_insights'] = []
    st.session_state['fetched_lp'] = 0
    st.session_state['fetched_zoning'] = ""
    st.session_state['rent_roll_data'] = [] 
    st.session_state['rent_roll_init'] = False

# --- [헬퍼 함수] ---
def get_address_from_coords(lat, lng):
    url = "http://api.vworld.kr/req/address" # HTTP 사용
    params = {"service": "address", "request": "getaddress", "version": "2.0", "crs": "EPSG:4326", "point": f"{lng},{lat}", "type": "PARCEL", "format": "json", "key": VWORLD_KEY}
    headers = {"Referer": REFERER_URL}
    try:
        res = requests.get(url, params=params, headers=headers, timeout=5, verify=False)
        if res.json().get('response', {}).get('status') == 'OK': return res.json()['response']['result'][0]['text']
    except: return None

def render_styled_block(label, value):
    st.markdown(f"<div style='margin-bottom:10px;'><div style='font-size:16px;color:#555;font-weight:700;'>{label}</div><div style='font-size:26px;font-weight:900;color:#111;'>{value}</div></div>", unsafe_allow_html=True)

def editable_area_input(label, key, default_val):
    val_str = st.text_input(label, value=str(default_val), key=key)
    try:
        val = float(str(val_str).replace(',', ''))
        st.markdown(f"<div style='color:#D32F2F;font-size:24px;font-weight:900;text-align:right;margin-top:-5px;'>{val*0.3025:,.1f} 평</div>", unsafe_allow_html=True)
        return val
    except: 
        st.markdown(f"<div style='color:#D32F2F;font-size:24px;font-weight:900;text-align:right;margin-top:-5px;'>- 평</div>", unsafe_allow_html=True)
        return 0.0

def editable_text_input(label, key, default_val):
    return st.text_input(label, value=str(default_val), key=key)

def comma_input(label, unit, key, default_val, help_text=""):
    st.markdown(f"<div style='font-size:17px;font-weight:800;color:#222;margin-bottom:4px;'>{label} <span style='font-size:13px;color:#666;'>{help_text}</span></div>", unsafe_allow_html=True)
    c1, c2 = st.columns([3, 1]) 
    with c1:
        if key not in st.session_state: st.session_state[key] = default_val
        val_in = st.text_input(label, value=f"{st.session_state[key]:,}" if st.session_state[key] else "", key=f"{key}_w", label_visibility="hidden")
        try: st.session_state[key] = int(str(val_in).replace(',', '').strip()) if val_in else 0
        except: st.session_state[key] = 0
    with c2: st.markdown(f"<div style='margin-top:15px;font-size:19px;font-weight:700;color:#444;'>{unit}</div>", unsafe_allow_html=True)
    return st.session_state[key]

def format_date_dot(date_str):
    if not date_str or len(date_str) != 8: return date_str
    return f"{date_str[:4]}.{date_str[4:6]}.{date_str[6:]}"

def generate_insight_candidates(info, finance, zoning, env_features, user_comment, comp_df=None, target_dong=""):
    points = []
    marketing_db = {
        "역세권": ["☑ [초역세권] 지하철역 도보 3분 이내, 유동인구 끊이지 않는 특급 입지", "☑ [교통허브] 대중교통 접근성 최상, 임차인 선호도 1순위 지역", "☑ [환금성] 역세권 프리미엄으로 경기 불황에도 시세 방어 탁월", "☑ [출퇴근용이] 직장인 수요 풍부하여 공실 위험 극히 낮은 안전 자산"],
        "광대로변": ["☑ [광대로] 왕복 8차선 이상 광대로 접함, 압도적인 웅장함 자랑", "☑ [랜드마크] 지역을 대표하는 대로변 건물로 기업 인지도 상승 효과", "☑ [접근성] 차량 진출입 및 대중교통 연계성 최고의 A급 입지", "☑ [상징성] 대기업 플래그십 스토어 및 전시장 용도로 강력 추천"],
        "먹자상권": ["☑ [먹자상권] 24시간 불 꺼지지 않는 상권, 권리금 형성된 검증된 자리", "☑ [유동인구] 점심부터 회식까지 직장인/거주민 발길 끊이지 않는 곳", "☑ [임대수요] 식당, 주점, 카페 등 임차 대기 수요 풍부한 핫플레이스", "☑ [현금흐름] 높은 바닥 권리금으로 공실 리스크 제로에 가까운 입지"],
        "랜드마크": ["☑ [랜드마크] 지역 내 누구나 아는 상징적인 건물, 투자가치 확실", "☑ [시세리딩] 주변 시세를 이끄는 대장주 건물, 자산 가치 상승 보장", "☑ [프리미엄] 압도적인 규모와 외관으로 지역 내 독보적인 존재감", "☑ [명품입지] 성공한 기업의 사옥이나 병원으로 최고의 선택"],
        "급매물": ["☑ [초급매] 시세 대비 20% 이상 저렴, 다시 없을 기회의 급매물", "☑ [안전마진] 매입 즉시 시세 차익 확정된 저평가 알짜 자산", "☑ [적극추천] 건물주 사정상 긴급 매각, 빠른 거래 시 네고 가능성", "☑ [투자찬스] 대지 평당가만 봐도 이득인 가격, 망설이면 놓치는 매물"],
        "주차편리": ["☑ [주차편리] 강남권에서 보기 드문 자주식 주차장 확보, 임차 경쟁력 甲", "☑ [대형차량] SUV, 고급 세단도 진입 가능한 넓은 주차 진입로 보유", "☑ [발렛파킹] 주차 부스 및 발렛 공간 확보로 방문객 편의성 극대화", "☑ [주차수익] 넉넉한 주차 공간 활용해 월 주차 추가 수익 창출 가능"],
        "감정가이하": ["☑ [감정가이하] 탁상감정가보다 낮은 매매가, 대출 한도 극대화 가능", "☑ [안전마진] 감정평가 금액보다 싸게 사는 확실한 안전 마진 확보", "☑ [LTV유리] 매가 대비 높은 대출 비율로 실투자금 최소화 전략", "☑ [투자1순위] 경매보다 싸게 살 수 있는 일반 매매 초급매물"],
        "초역세권": ["☑ [초역세권] 지하철 출구 나오자마자 보이는 건물, 접근성 끝판왕", "☑ [유동인구] 출퇴근길 필수 동선, 자연스러운 워크인 고객 유입", "☑ [희소성] 역 바로 앞 토지는 부르는 게 값, 소장 가치 100%", "☑ [광고효과] 지하철 이용객에게 24시간 노출되는 옥외 광고 명당"],
        "대로변": ["☑ [대로변] 가시성 최상급, 기업 홍보 효과 탁월한 랜드마크 입지", "☑ [Trophy Asset] 소장 가치 확실한 대로변 빌딩, 자산 가치 상승 기대", "☑ [상징성] 웅장한 전면 확보로 대형 프랜차이즈 및 병원 임차 유리", "☑ [안전성] 환금성 뛰어난 대로변 입지, 언제든 현금화 가능"],
        "오피스상권": ["☑ [오피스상권] 구매력 높은 직장인 상주, 평일 점심 매출 폭발적", "☑ [B2B수요] 주변 대기업 협력사 및 관련 업체 사무실 수요 풍부", "☑ [인프라] 은행, 관공서 등 업무 지원 시설 인접해 사옥으로 안성맞춤", "☑ [주5일상권] 주말 관리 용이하고 평일 집중적인 매출 발생하는 효율적 상권"],
        "법조타운": ["☑ [법조타운] 법원/검찰청 인접, 변호사 및 법무사 사무실 임차 수요 풍부", "☑ [전문직] 월세 밀릴 걱정 없는 고소득 전문직 임차인 선호 지역", "☑ [안정성] 경기 타지 않는 법률 관련 업종 밀집, 꾸준한 임대 수익", "☑ [희소입지] 법조타운 내 건물은 매물이 귀해 나오면 바로 거래됨"],
        "사옥추천": ["☑ [사옥추천] 쾌적한 업무 공간과 넉넉한 주차, 임직원 만족도 최상", "☑ [브랜딩] 세련된 외관으로 기업 이미지 상승 및 홍보 효과 기대", "☑ [효율성] 전용률 높고 레이아웃 배치 유리해 공간 활용도 200%", "☑ [자산가치] 사옥으로 사용하며 지가 상승 누리는 '일석이조' 투자"],
        "수려한외관": ["☑ [디자인] 건축상 수상급의 수려한 외관, 지역 내 시선 집중", "☑ [컨디션] 내외관 관리 상태 최상, 손볼 곳 없이 즉시 수익 발생", "☑ [가시성] 독특하고 아름다운 외관으로 자연스러운 홍보 효과", "☑ [임차유리] 디자인 에이전시, IT 기업 등 감각적인 임차인 선호"],
        "용적률이득": ["☑ [용적률이득] 현재 법정 용적률보다 더 높게 지어진 '이득 본' 건물", "☑ [가성비] 같은 땅 크기라도 연면적이 넓어 임대 수익 극대화 가능", "☑ [희소성] 지금 신축하면 이만큼 못 짓는 귀한 '오버 용적률' 매물", "☑ [철거반대] 신축보다는 리모델링을 통해 기존 면적 살리는 전략 추천"],
        "더블역세권": ["☑ [더블역세권] 2개 노선이 교차하는 교통 요충지, 광역 수요 흡수", "☑ [황금노선] 강남 및 주요 업무지구 접근성 우수, 지가 상승 견인", "☑ [희소성] 더블 역세권의 희소가치로 향후 매각 차익 극대화 기대", "☑ [S급입지] 유동인구와 배후수요 모두 갖춘 실패 없는 투자처"],
        "대로코너": ["☑ [대로코너] 대로변과 이면도로를 모두 접한 최고의 명당 자리", "☑ [가시성] 3면 개방형으로 어디서든 잘 보이는 압도적 노출 효과", "☑ [진입성] 차량 및 보행자 접근이 모두 용이한 최적의 입지", "☑ [프리미엄] 코너 자리는 부르는 게 값, 향후 지가 상승 1순위"],
        "항아리상권": ["☑ [독점상권] 5,000세대 이상 대단지 배후 수요 독점하는 항아리 입지", "☑ [생활밀착] 병원, 약국, 학원 등 필수 근생 업종 최적화된 건물", "☑ [충성고객] 외부 유출 없이 단지 내 고정 고객 확보된 안정적 상권", "☑ [공실제로] 한 번 들어오면 나가지 않는 임차인, 관리 편한 효자 매물"],
        "핫플레이스": ["☑ [핫플레이스] MZ세대가 찾아오는 지역 명소, SNS 업로드 성지", "☑ [트렌드] 가장 힙한 F&B 브랜드들이 입점하고 싶어 하는 건물", "☑ [미래가치] 상권 확장세가 뚜렷하여 매각 시 높은 시세 차익 기대", "☑ [권리금] 높은 권리금이 형성되어 있어 임대료 연체 걱정 없음"],
        "수익형": ["☑ [수익형] 탄탄한 임차 구성으로 매월 따박따박 월세 들어오는 효자", "☑ [고수익] 레버리지 활용 시 연 5% 이상 수익률 달성 가능한 알짜", "☑ [공실제로] 대기 수요 풍부해 공실 걱정 없이 마음 편한 임대 사업", "☑ [은퇴준비] 안정적인 현금 흐름으로 노후 대비 및 증여용으로 강력 추천"],
        "신축빌딩": ["☑ [랜드마크] 지역 내 압도적인 존재감 자랑하는 최신축 하이엔드 빌딩", "☑ [희소성] 노후 건물 즐비한 지역 내 단비 같은 신축, 경쟁 우위 확실", "☑ [프리미엄] 신축 프리미엄으로 향후 매각 시 높은 시세 차익 기대", "☑ [절세효과] 법인 사옥 매입 시 비용 처리 및 자산 가치 상승 동시 효과"],
        "신축부지용": ["☑ [신축부지] 노후 건물을 철거하고 원하는 컨셉으로 신축 가능한 땅", "☑ [디벨로핑] 명도 완료 후 신축 시 드라마틱한 가치 상승 기대", "☑ [맞춤설계] 사옥이나 메디컬 빌딩 등 내 입맛대로 건축 가능", "☑ [토지가치] 건물값은 0원, 오직 땅의 가치만 보고 투자하는 곳"],
        "트리플역세권": ["☑ [트리플역세권] 3개 노선이 만나는 교통의 심장부, 폭발적 유동인구", "☑ [초특급입지] 서울 전역 어디든 빠르게 이동 가능한 교통 허브", "☑ [투자가치] 대한민국 상위 1% 입지, 묻어두면 무조건 오르는 곳", "☑ [임대보장] 공실이란 단어가 없는 곳, 최고의 임대 안정성 자랑"],
        "이면코너": ["☑ [이면코너] 메인 상권 이면의 알짜 코너, 가성비 최고의 실속 매물", "☑ [먹자상권] 점심/저녁 유동인구 바글바글한 이면 먹자 골목의 요지", "☑ [안정성] 유행 타지 않는 탄탄한 배후 수요로 공실 걱정 없는 곳", "☑ [수익률] 대로변 대비 합리적 매가로 높은 임대 수익 실현 가능"],
        "학군지": ["☑ [학군지] 대치/목동급 명문 학원가 밀집, 학원 임차 수요 폭발", "☑ [항아리상권] 학생 및 학부모 고정 수요로 365일 북적이는 상권", "☑ [우량임차] 프랜차이즈 학원, 스터디카페 등 장기 우량 임차인 선호", "☑ [불패신화] 교육열 높은 지역 특성상 경기 침체에도 끄떡없는 상권"],
        "메디컬입지": ["☑ [메디컬] 엘리베이터, 주차, 전력 등 병의원 개원 하드웨어 완벽", "☑ [독점수요] 약국 입점 가능해 고수익 창출 및 건물 가치 레벨업", "☑ [시너지] 내과, 이비인후과 등 연계 처방 가능한 메디컬 빌딩 최적", "☑ [고령화] 인구 고령화로 의료 수요 증가, 공실 걱정 없는 메디컬 입지"],
        "시세차익": ["☑ [시세차익] 지금 사두면 3년 뒤 웃게 될 확실한 지가 상승 지역", "☑ [저평가] 주변 호재 대비 아직 저평가된 가격, 상승 여력 충분", "☑ [개발호재] GTX, 재개발 등 대형 호재가 대기 중인 투자 유망처", "☑ [토지투자] 건물 수익보다 땅값 상승으로 자산을 불릴 최고의 기회"],
        "관리상태최상": ["☑ [관리최상] 건물주가 거주하며 직접 관리해 신축처럼 깔끔한 컨디션", "☑ [비용절감] 누수/방수/외벽 등 최근 보수 완료, 손볼 곳 하나 없음", "☑ [임차인] 건물 관리 잘되어 임차인 만족도 높고 재계약률 우수", "☑ [가성비] 연식 대비 너무나 깨끗한 내외관, 가성비 최고의 선택"],
        "명도완료": ["☑ [명도완료] 골치 아픈 명도 문제 100% 해결, 즉시 사업 진행 가능", "☑ [신축추천] 잔금과 동시에 철거 가능하여 금융 비용 절감 효과", "☑ [리모델링] 텅 빈 건물에서 내가 원하는 대로 대수선 공사 가능", "☑ [시간절약] 명도 협의에 걸리는 수개월의 시간을 번 셈인 귀한 매물"],
        "광역환승": ["☑ [광역환승] GTX/KTX/광역버스 환승 센터 인접, 교통의 집결지", "☑ [미래가치] 광역 교통망 확충으로 향후 천지개벽할 투자의 중심", "☑ [유동인구] 경기도 및 지방 인구까지 흡수하는 거대 상권 형성", "☑ [개발호재] 환승센터 복합 개발 등 대형 호재의 직접 수혜지"],
        "이면초입": ["☑ [이면초입] 대로변 바로 한 블록 뒤, 가시성과 가성비 모두 잡음", "☑ [진입로] 상권으로 들어가는 관문 입지, 유동인구 자연 유입", "☑ [실속형] 대로변의 광고 효과는 누리면서 가격은 합리적인 매물", "☑ [먹자초입] 먹자골목 시작점에 위치하여 만남의 장소로 유명"],
        "숲세권": ["☑ [숲세권] 도심 속 힐링, 공원 및 녹지 인접해 쾌적한 업무 환경", "☑ [워라밸] 점심시간 산책 가능한 숲세권 입지, 젊은 직장인 선호도 급상승", "☑ [뷰맛집] 창밖으로 펼쳐지는 파크뷰, 개방감과 조망권 확보된 건물", "☑ [희소성] 팍팍한 도심 내 희소한 자연 친화적 입지, 미래 가치 우수"],
        "가시성우수": ["☑ [가시성] 멀리서도 한눈에 들어오는 탁월한 가시성, 광고판 그 자체", "☑ [노출효과] 차량 및 보행자에게 24시간 자연스럽게 노출되는 위치", "☑ [홍보명당] 간판 설치 시 광고 효과가 뛰어나 임차인이 선호하는 곳", "☑ [개방감] 답답함 없이 뻥 뚫린 시야 확보, 건물이 커 보이는 효과"],
        "벨류업유망": ["☑ [밸류업] 리모델링 시 임대료 2배 상승 가능한 원석 같은 매물", "☑ [용적률] 법적 용적률 대비 덜 찾아먹은 상태, 증축 통해 가치 극대화", "☑ [디벨로퍼] 낡은 건물을 명도 후 신축하면 수익률 30% 이상 기대", "☑ [가치투자] 현재의 낡은 모습보다 미래의 바뀐 모습을 보고 투자할 곳"],
        "리모델링": ["☑ [리모델링] 뼈대만 남기고 싹 고치면 신축급 성능 발휘 가능", "☑ [비용절감] 신축 대비 저렴한 공사비로 드라마틱한 임대료 상승", "☑ [트렌드] 노후 건물을 힙한 공간으로 변신시켜 젊은 층 유입 유도", "☑ [엘리베이터] 승강기 신설 및 외관 교체 시 건물 가치 수직 상승"],
        "명도협의가능": ["☑ [명도협의] 매수인의 계획에 맞춰 명도 시기 및 조건 유연하게 조율", "☑ [책임명도] 매도인이 책임지고 명도해주는 조건, 골치 아플 일 없음", "☑ [실사용] 일부 층 명도하여 매수인이 즉시 입주 가능한 실용적 매물", "☑ [윈윈] 기존 임차인 승계와 명도 중 유리한 쪽으로 선택 가능"]
    }
    
    final_results = []
    if user_comment: final_results.append(f"📌 {user_comment.strip()[:40]}") 
    num_selected = len(env_features)
    target_count = 10
    if num_selected > 0:
        if num_selected <= target_count:
            base_count = target_count // num_selected
            remainder = target_count % num_selected
            for i, keyword in enumerate(env_features):
                if keyword in marketing_db:
                    count_to_pick = base_count + (1 if i < remainder else 0)
                    pool = marketing_db[keyword]
                    real_count = min(len(pool), count_to_pick)
                    picked = random.sample(pool, real_count)
                    final_results.extend(picked)
        else:
            chosen_keywords = random.sample(env_features, target_count)
            for keyword in chosen_keywords:
                if keyword in marketing_db:
                    final_results.extend(random.sample(marketing_db[keyword], 1))
    unique_final_points = list(dict.fromkeys(final_results))
    if comp_df is not None and not comp_df.empty:
        try:
            sold_df = comp_df[comp_df['구분'].astype(str).str.contains('매각|완료|매매', na=False)]
            if not sold_df.empty:
                avg_price = sold_df['평당가'].mean()
                my_price = finance['land_pyeong_price_val']
                diff = my_price - avg_price
                diff_pct = abs(diff / avg_price) * 100
                loc_text = target_dong if target_dong else "인근"
                if diff < 0: points.append(random.choice([f"☑ [가격우위] {loc_text} 평균(평 {avg_price:,.0f}만) 대비 {diff_pct:.1f}% 저렴한 저평가 매물", f"☑ [가격메리트] 주변 시세보다 평당 {abs(diff):,.0f}만원 싸게 나온 확실한 급매"]))
                else: points.append(random.choice([f"☑ [대장주] {loc_text} 시세를 리딩하는 압도적 컨디션의 대장 건물", f"☑ [프리미엄] 평균보다 높지만 그만한 가치가 있는 A급 입지"]))
        except: pass
    yield_val = finance['yield']
    if yield_val >= 4.5: points.append(random.choice([f"☑ [초고수익] 연 {yield_val:.1f}% 수익률! 요즘 같은 고금리에 보기 드문 보물", f"☑ [현금흐름] 묻어두면 돈이 되는 연 {yield_val:.1f}% 수익형 부동산 끝판왕"]))
    elif yield_val >= 3.5: points.append(random.choice([f"☑ [고수익] 연 {yield_val:.1f}% 안정적 수익, 이자 내고도 남는 훌륭한 수익성", f"☑ [알짜매물] 수익률과 지가 상승 두 마리 토끼 잡는 연 {yield_val:.1f}% 매물"]))
    elif yield_val >= 2.5: points.append(random.choice([f"☑ [안정성] 연 {yield_val:.1f}%의 꾸준한 임대 수익과 확실한 지가 상승 동시 추구", f"☑ [리스크헷지] 공실 걱정 없는 입지에서 누리는 연 {yield_val:.1f}%의 편안함"]))
    else: points.append(random.choice([f"☑ [미래가치] 당장 수익보다 향후 폭발적 지가 상승과 개발 호재에 집중", f"☑ [시세차익] 보유할수록 땅값이 오르는 토지 가치 중심의 투자처"]))
    fallback_msgs = ["☑ [희소가치] 매물 잠김 심한 이 지역 내 오랜만에 등장한 귀한 물건", "☑ [육각형] 입지, 가격, 상권, 미래가치 4박자 모두 갖춘 보기 드문 투자처", "☑ [불패입지] 한번 들어오면 나가지 않는 임차인 선호도 1위 검증된 자리"]
    random.shuffle(fallback_msgs)
    points.extend(fallback_msgs)
    unique_final_points = list(dict.fromkeys(unique_final_points + points))
    return unique_final_points[:10]

# --- [API 조회] 핵심 수정 부분: HTTP 사용 ---
@st.cache_data(show_spinner=False)
def get_pnu_and_coords(address):
    # [수정] https -> http (502 에러 방지)
    url = "http://api.vworld.kr/req/search" 
    params = {
        "service": "search", "request": "search", "version": "2.0", 
        "crs": "EPSG:4326", "size": "1", "page": "1", 
        "query": address, "type": "address", "category": "parcel", 
        "format": "json", "errorformat": "json", "key": VWORLD_KEY
    }
    # [수정] 단순 헤더
    headers = {"Referer": REFERER_URL}
    try:
        res = requests.get(url, params=params, headers=headers, timeout=10) # 타임아웃 10초로 늘림
        if res.status_code == 200:
            data = res.json()
            if data.get('response', {}).get('status') == 'OK': 
                item = data['response']['result']['items'][0]
                pnu = item.get('address', {}).get('pnu') or item.get('id')
                lng = float(item['point']['x']); lat = float(item['point']['y'])
                full_address = item.get('address', {}).get('parcel', '') or item.get('address', {}).get('road', '') or address
                return {"pnu": pnu, "lat": lat, "lng": lng, "full_addr": full_address}
            return None
        else:
            st.error(f"서버 응답 코드: {res.status_code}")
            return None
    except Exception as e:
        st.error(f"연결 오류: {e}")
        return None

@st.cache_data(show_spinner=False)
def get_zoning_smart(lat, lng):
    # [수정] https -> http
    url = "http://api.vworld.kr/req/data"
    delta = 0.0005
    min_x, min_y = lng - delta, lat - delta
    max_x, max_y = lng + delta, lat + delta
    params = {"service": "data", "request": "GetFeature", "data": "LT_C_UQ111", "key": VWORLD_KEY, "format": "json", "size": "10", "geomFilter": f"BOX({min_x},{min_y},{max_x},{max_y})", "domain": REFERER_URL}
    headers = {"Referer": REFERER_URL}
    try:
        res = requests.get(url, params=params, headers=headers, timeout=5)
        if res.status_code == 200:
            features = res.json().get('response', {}).get('result', {}).get('featureCollection', {}).get('features', [])
            if features: return ", ".join(sorted(list(set([f['properties']['UNAME'] for f in features]))))
    except: pass
    return ""

@st.cache_data(show_spinner=False)
def get_land_price(pnu):
    url = "https://apis.data.go.kr/1611000/NsdiIndvdLandPriceService/getIndvdLandPriceAttr"
    current_year = datetime.datetime.now().year
    for year in range(current_year, current_year - 7, -1):
        params = {"serviceKey": USER_KEY, "pnu": pnu, "format": "xml", "numOfRows": "1", "pageNo": "1", "stdrYear": str(year)}
        try:
            res = requests.get(url, params=params, timeout=4)
            if res.status_code == 200:
                root = ET.fromstring(res.content)
                if root.findtext('.//resultCode') == '00':
                    price_node = root.find('.//indvdLandPrice')
                    if price_node is not None and price_node.text: return int(price_node.text)
        except: continue
        time.sleep(0.05)
    return 0

@st.cache_data(show_spinner=False)
def get_building_info_smart(pnu):
    base_url = "https://apis.data.go.kr/1613000/BldRgstHubService/getBrTitleInfo"
    sigungu = pnu[0:5]; bjdong = pnu[5:10]; bun = pnu[11:15]; ji = pnu[15:19]
    plat_code = '1' if pnu[10] == '2' else '0'
    params = {"serviceKey": USER_KEY, "sigunguCd": sigungu, "bjdongCd": bjdong, "platGbCd": plat_code, "bun": bun, "ji": ji, "numOfRows": "1", "pageNo": "1"}
    try:
        res = requests.get(base_url, params=params, timeout=5, verify=False)
        if res.status_code == 200: return parse_xml_response(res.content)
        return {"error": f"서버 상태: {res.status_code}"}
    except Exception as e: return {"error": str(e)}

def parse_xml_response(content):
    try:
        root = ET.fromstring(content)
        item = root.find('.//item')
        if item is None: return None
        indr_mech = int(item.findtext('indrMechUtcnt', '0') or 0); indr_auto = int(item.findtext('indrAutoUtcnt', '0') or 0)
        oudr_mech = int(item.findtext('oudrMechUtcnt', '0') or 0); oudr_auto = int(item.findtext('oudrAutoUtcnt', '0') or 0)
        ride_elvt = int(item.findtext('rideUseElvtCnt', '0') or 0); emgen_elvt = int(item.findtext('emgenUseElvtCnt', '0') or 0)
        return {
            "bldNm": item.findtext('bldNm', '-'), "mainPurpsCdNm": item.findtext('mainPurpsCdNm', '정보없음'),
            "strctCdNm": item.findtext('strctCdNm', '정보없음'), "platArea": float(item.findtext('platArea', '0') or 0),
            "totArea": float(item.findtext('totArea', '0') or 0), "archArea_val": float(item.findtext('archArea', '0') or 0),
            "groundArea": float(item.findtext('vlRatEstmTotArea', '0') or 0), "ugrndFlrCnt": int(item.findtext('ugrndFlrCnt', '0') or 0),
            "grndFlrCnt": int(item.findtext('grndFlrCnt', '0') or 0), "useAprDay": format_date_dot(item.findtext('useAprDay', '')),
            "bcRat": float(item.findtext('bcRat', '0') or 0), "vlRat": float(item.findtext('vlRat', '0') or 0),
            "rideUseElvtCnt": f"{ride_elvt + emgen_elvt}대", "parking": f"{indr_mech+indr_auto+oudr_mech+oudr_auto}대(옥내{indr_mech+indr_auto}/옥외{oudr_mech+oudr_auto})"
        }
    except Exception as e: return {"error": str(e)}

@st.cache_data(show_spinner=False)
def get_floor_info_smart(pnu):
    base_url = "https://apis.data.go.kr/1613000/BldRgstHubService/getBrFlrOulnInfo"
    sigungu = pnu[0:5]; bjdong = pnu[5:10]; bun = pnu[11:15]; ji = pnu[15:19]
    plat_code = '1' if pnu[10] == '2' else '0'
    params = {"serviceKey": USER_KEY, "sigunguCd": sigungu, "bjdongCd": bjdong, "platGbCd": plat_code, "bun": bun, "ji": ji, "numOfRows": "50", "pageNo": "1"}
    floor_data = []
    try:
        res = requests.get(base_url, params=params, timeout=5)
        if res.status_code == 200:
            root = ET.fromstring(res.content)
            items = root.findall('.//item')
            floor_map = {} 
            for item in items:
                try:
                    flr_no = int(item.findtext('flrNo')); flr_gb = item.findtext('flrGbCdNm')
                    area = float(item.findtext('area', '0') or 0)
                    idx = -flr_no if '지하' in flr_gb else flr_no
                    floor_map[idx] = floor_map.get(idx, 0) + area
                except: continue
            for idx in sorted(floor_map.keys()):
                flr_name = f"B{abs(idx)}" if idx < 0 else f"{idx}층"
                floor_data.append({"층수": flr_name, "입주업체": "", "층별면적": f"{floor_map[idx]*0.3025:.2f}", "보증금": None, "임대료": None, "관리비": None, "임대차기간": "", "비고": ""})
            return floor_data
    except Exception as e: print(e)
    return []

@st.cache_data(show_spinner=False)
def get_cadastral_map_image(lat, lng):
    delta = 0.0015 
    bbox = f"{lng-delta},{lat-delta},{lng+delta},{lat+delta}"
    # [수정] https -> http
    url = f"http://api.vworld.kr/req/wms?SERVICE=WMS&REQUEST=GetMap&VERSION=1.3.0&LAYERS=LP_PA_CBND_BUBUN&STYLES=LP_PA_CBND_BUBUN&CRS=EPSG:4326&BBOX={bbox}&WIDTH=400&HEIGHT=300&FORMAT=image/png&TRANSPARENT=FALSE&BGCOLOR=0xFFFFFF&EXCEPTIONS=text/xml&KEY={VWORLD_KEY}"
    headers = {"Referer": REFERER_URL}
    try:
        res = requests.get(url, headers=headers, timeout=5)
        if res.status_code == 200 and 'image' in res.headers.get('Content-Type', ''): return BytesIO(res.content)
    except: pass
    return None

@st.cache_data(show_spinner=False)
def get_static_map_image(lat, lng):
    # [수정] https -> http
    url = f"http://api.vworld.kr/req/image?service=image&request=getmap&key={VWORLD_KEY}&center={lng},{lat}&crs=EPSG:4326&zoom=17&size=600,400&format=png&basemap=GRAPHIC"
    headers = {"Referer": REFERER_URL}
    try:
        res = requests.get(url, headers=headers, timeout=3)
        if res.status_code == 200 and 'image' in res.headers.get('Content-Type', ''): return BytesIO(res.content)
    except: pass
    return None

def create_pptx(info, full_addr, finance, zoning, lat, lng, land_price, selling_points, images_dict, rent_roll_df=None, template_binary=None, template_1page_binary=None):
    deep_red = RGBColor(204, 0, 0); black = RGBColor(0, 0, 0)
    bld_name = info.get('bldNm') if info.get('bldNm') != '-' else f"{full_addr.split(' ')[2] if len(full_addr.split(' ')) > 2 else ''} 빌딩"
    lp_str_final = f"{(land_price/10000)/0.3025:,.0f}만원/평" if land_price > 0 else "0만원/평"
    total_lp_num = int((land_price * info['platArea']) / 100000000) if land_price and info['platArea'] else 0
    
    ctx_vals = {'plat_m2': f"{info['platArea']:,}", 'plat_py': f"{info['platArea'] * 0.3025:,.1f}", 
                'tot_m2': f"{info['totArea']:,}", 'tot_py': f"{info['totArea'] * 0.3025:,.1f}", 
                'arch_m2': f"{info.get('archArea_val', 0):,.1f}", 'arch_py': f"{info.get('archArea_val', 0) * 0.3025:,.1f}", 
                'ground_m2': f"{info.get('groundArea', 0):,}", 'ground_py': f"{info.get('groundArea', 0) * 0.3025:,.1f}", 
                'use_date': info.get('useAprDay', '-')}

    data_map = {
        "{{빌딩이름}}": bld_name, "{{소재지}}": full_addr, "{{용도지역}}": zoning,
        "{{AI물건분석내용 4가지 }}": "\n".join(selling_points[:5]) if selling_points else "분석된 특징이 없습니다.", 
        "{{공시지가}}": lp_str_final, "{{공시지가 총액}}": f"합 {total_lp_num:,}억" if total_lp_num > 0 else "-",
        "{{교통편의}}": info.get('traffic', '-'), "{{도로상황}}": info.get('road', '-'), "{{준공년도}}": ctx_vals['use_date'],
        "{{건물규모}}": info.get('scale_str', '-'), "{{건폐율}}": info.get('bcRat_str', '-'), "{{용적률}}": info.get('vlRat_str', '-'), 
        "{{승강기}}": info.get('rideUseElvtCnt', '-'), "{{주차대수}}": info.get('parking', '-'), "{{건물주구조}}": info.get('strctCdNm', '-'),
        "{{건물용도}}": info.get('mainPurpsCdNm', '-'), "{{보증금}}": f"{finance['deposit']:,} 만원", "{{월임대료}}": f"{finance['rent']:,} 만원",
        "{{관리비}}": f"{finance['maintenance']:,} 만원", "{{수익률}}": f"년 {finance['yield']:.1f}%", "{{융자금}}": f"{finance['loan']:,} 억원",
        "{{매매금액}}": f"{finance['price']:,} 억원", "{{대지평단가}}": f"평 {finance.get('land_pyeong_price_val', 0):,.0f}만원",
        "{{건물미래가치 활용도}}": "사옥 및 수익용 리모델링 추천", "{{위치도}}": "", "{{지적도}}": "", "{{건축물대장}}": "", "{{건물사진}}": ""
    }

    def replace_text_in_frame(text_frame, mapper, ctx):
        for p in text_frame.paragraphs:
            p_text = p.text
            if "{{매매금액}}" in p_text or "{{수익률}}" in p_text or "{{공시지가 총액}}" in p_text:
                key = [k for k in ["{{매매금액}}", "{{수익률}}", "{{공시지가 총액}}"] if k in p_text][0]
                p.text = str(mapper[key])
                for r in p.runs: r.font.bold = True; r.font.color.rgb = deep_red; r.font.size = Pt(16) if key == "{{매매금액}}" else Pt(12)
                return
            for k, v in mapper.items():
                if k in p_text and k not in ["{{매매금액}}", "{{수익률}}", "{{공시지가 총액}}"]:
                    p.text = p_text.replace(k, str(v)); return
            if "{{대지면적}}" in p_text:
                p.text = p_text.replace("{{대지면적}}", ctx['plat_py'] if "평" in p_text else ctx['plat_m2'])
                for r in p.runs: r.font.bold = True; r.font.color.rgb = RGBColor(6, 6, 236) if "평" in p_text else black
            elif "{{연면적}}" in p_text:
                p.text = p_text.replace("{{연면적}}", ctx['tot_py'] if "평" in p_text else ctx['tot_m2'])
                for r in p.runs: r.font.bold = True; r.font.color.rgb = RGBColor(6, 6, 236) if "평" in p_text else black
            elif "{{건축면적}}" in p_text: p.text = p_text.replace("{{건축면적}}", ctx['arch_py'] if "평" in p_text else ctx['arch_m2'])
            elif "{{지상면적}}" in p_text: p.text = p_text.replace("{{지상면적}}", ctx['ground_py'] if "평" in p_text else ctx['ground_m2'])
            elif "{{준공년도}}" in p_text: p.text = p_text.replace("{{준공년도}}", ctx['use_date'])

    def replace_text_in_shape(shape, mapper, ctx):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes: replace_text_in_shape(child, mapper, ctx)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame: replace_text_in_frame(cell.text_frame, mapper, ctx)
        elif shape.has_text_frame: replace_text_in_frame(shape.text_frame, mapper, ctx)

    if template_binary:
        prs = Presentation(template_binary)
        for slide in prs.slides:
            for shape in slide.shapes: replace_text_in_shape(shape, data_map, ctx_vals)
        img_insert_map = {1: 'u1', 2: 'u2', 4: 'u3', 5: 'u4'}
        for s_idx, key in img_insert_map.items():
            if s_idx < len(prs.slides) and images_dict.get(key):
                f = images_dict[key]; f.seek(0)
                prs.slides[s_idx].shapes.add_picture(f, Cm(1.35), Cm(2.35), width=Cm(24.84) if key=='u1' else Cm(13.61) if key=='u2' else Cm(20.4), height=Cm(15.74))
        if 6 < len(prs.slides):
            for i, k in enumerate(['u5_1', 'u5_2', 'u5_3', 'u5_4']):
                if images_dict.get(k):
                    f = images_dict[k]; f.seek(0)
                    prs.slides[6].shapes.add_picture(f, Cm(1.35 + (12.48 * (i%2))), Cm(2.35 + (8.2 * (i//2))), width=Cm(12.16), height=Cm(7.74))
        if rent_roll_df is not None and not rent_roll_df.empty:
            try:
                tbl = [s for s in prs.slides[3].shapes if s.has_table][0].table
                for i, row in enumerate(rent_roll_df.to_dict('records')):
                    if i + 1 >= len(tbl.rows) - 1: break
                    for j, col in enumerate(['층수', '입주업체', '층별면적', '보증금', '임대료', '관리비', '임대차기간', '비고']):
                         if j < len(tbl.columns): tbl.cell(i+1, j+1).text = str(row.get(col, '') or '')
            except: pass
        out = BytesIO(); prs.save(out); return out.getvalue()
    elif template_1page_binary:
        prs = Presentation(template_1page_binary)
        img_specs = {"{{건물사진}}": (8.78, 11.11), "{{위치도}}": (8.78, 9.17), "{{지적도}}": (9.03, 5.9)}
        for slide in prs.slides:
            for shape in list(slide.shapes):
                if shape.has_text_frame and shape.text_frame.text.strip() in img_specs:
                    k = shape.text_frame.text.strip()
                    img_data = images_dict.get('u2' if k=="{{건물사진}}" else 'u1' if k=="{{위치도}}" else 'u3') or (get_static_map_image(lat, lng) if k=="{{위치도}}" else get_cadastral_map_image(lat, lng) if k=="{{지적도}}" else None)
                    if img_data:
                        img_data.seek(0)
                        slide.shapes.add_picture(img_data, shape.left, shape.top, width=Cm(img_specs[k][0]), height=Cm(img_specs[k][1]))
                        shape._element.getparent().remove(shape._element)
            for shape in slide.shapes: replace_text_in_shape(shape, data_map, ctx_vals)
        out = BytesIO(); prs.save(out); return out.getvalue()
    else:
        prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Cm(1), Cm(1), Cm(19), Cm(2)).text_frame.text = bld_name
        out = BytesIO(); prs.save(out); return out.getvalue()

def create_excel(info, full_addr, finance, zoning, lat, lng, land_price, selling_points, uploaded_img):
    output = BytesIO(); workbook = xlsxwriter.Workbook(output, {'in_memory': True}); worksheet = workbook.add_worksheet('부동산분석')
    fmt = workbook.add_format({'border': 1, 'align': 'center', 'valign': 'vcenter'})
    worksheet.merge_range('B2:J3', info.get('bldNm', '-'), fmt)
    if uploaded_img: uploaded_img.seek(0); worksheet.insert_image('B6', 'img.png', {'image_data': uploaded_img, 'x_scale': 0.5, 'y_scale': 0.5})
    workbook.close(); return output.getvalue()

# =========================================================
# [메인 실행]
# =========================================================
st.title("🏢 부동산 매입 분석기 Pro")
st.markdown("---")

with st.expander("🗺 지도에서 직접 클릭하여 찾기 (Click)", expanded=False):
    m = folium.Map(location=[37.5172, 127.0473], zoom_start=14)
    output = st_folium(m, width=700, height=400)
    if output and output.get("last_clicked"):
        lat = output["last_clicked"]["lat"]; lng = output["last_clicked"]["lng"]
        if "last_click_lat" not in st.session_state or st.session_state["last_click_lat"] != lat:
            st.session_state["last_click_lat"] = lat
            found_addr = get_address_from_coords(lat, lng)
            if found_addr:
                st.success(f"📍 지도 클릭 확인! 변환된 주소: {found_addr}")
                st.session_state['addr'] = found_addr; reset_analysis(); st.rerun()
            else: st.warning("⚠️ 주소를 찾을 수 없는 위치입니다.")

link_container = st.container()
addr_input = st.text_input("주소 입력", placeholder="예: 강남구 논현동 254-4", key="addr", on_change=reset_analysis)

if addr_input:
    with st.spinner("데이터 분석 중..."):
        location = get_pnu_and_coords(addr_input)
        if not location: st.warning("⚠️ 해당주소가 없습니다")
        else:
            with link_container:
                c1, c2 = st.columns(2)
                c1.markdown(f"<a href='https://map.naver.com/v5/search/{quote_plus(location['full_addr'])}' target='_blank' class='link-btn naver-btn'>📍 네이버지도 위치확인</a>", unsafe_allow_html=True)
                if location.get('pnu'): c2.markdown(f"<a href='https://www.eum.go.kr/web/ar/lu/luLandDet.jsp?pnu={location['pnu']}&mode=search&isNoScr=script' target='_blank' class='link-btn eum-btn'>📑 토지이음 규제정보 확인</a>", unsafe_allow_html=True)
            
            if not st.session_state['zoning']: st.session_state['zoning'] = get_zoning_smart(location['lat'], location['lng'])
            if not st.session_state['fetched_zoning']: st.session_state['fetched_zoning'] = st.session_state['zoning']
            info = get_building_info_smart(location['pnu'])
            land_price = get_land_price(location['pnu'])
            if land_price > 0 and st.session_state['fetched_lp'] == 0: st.session_state['fetched_lp'] = land_price
            
            if not info or "error" in info: st.error(f"조회 실패: {info.get('error')}")
            else:
                st.success("✅ 분석 완료!")
                st.write("##### 📸 PPT 삽입용 사진 업로드")
                c1, c2, c3, c4 = st.columns(4)
                images_map = {
                    'u1': c1.file_uploader("Slide 2: 위치도", type=['png', 'jpg'], key="u1"),
                    'u2': c2.file_uploader("Slide 3: 건물메인", type=['png', 'jpg'], key="u2"),
                    'u3': c3.file_uploader("Slide 5: 지적도", type=['png', 'jpg'], key="u3"),
                    'u4': c4.file_uploader("Slide 6: 대장", type=['png', 'jpg'], key="u4")
                }
                st.write("▼ 추가 사진 (Slide 7)")
                cc1, cc2, cc3, cc4 = st.columns(4)
                images_map.update({'u5_1': cc1.file_uploader("추가1", key="u5_1"), 'u5_2': cc2.file_uploader("추가2", key="u5_2"), 'u5_3': cc3.file_uploader("추가3", key="u5_3"), 'u5_4': cc4.file_uploader("추가4", key="u5_4")})

                st.markdown("---")
                
                st.subheader("데이터 확인 및 수정")
                c_1, c_2 = st.columns([2, 1])
                c_1.text_input("소재지", value=addr_input, disabled=True)
                info['bldNm'] = c_2.text_input("건물명", value=info.get('bldNm', '-'))
                
                c_3, c_4, c_5 = st.columns(3)
                land_price = int(c_3.text_input("공시지가(원/㎡)", value=f"{st.session_state['fetched_lp']:,}").replace(',', ''))
                info['platArea'] = float(c_4.text_input("대지면적(㎡)", value=info['platArea']).replace(',', ''))
                info['totArea'] = float(c_5.text_input("연면적(㎡)", value=info['totArea']).replace(',', ''))
                
                st.subheader("📋 층별 임대 현황 (Rent Roll)")
                if not st.session_state['rent_roll_init']:
                     st.session_state['rent_roll_data'] = get_floor_info_smart(location['pnu']) or [{"층수": "1층", "입주업체": "", "층별면적": "", "보증금": 0, "임대료": 0, "관리비": 0, "임대차기간": "", "비고": ""}]
                     st.session_state['rent_roll_init'] = True
                
                edited_df = st.data_editor(st.session_state['rent_roll_data'], num_rows="dynamic", use_container_width=True, key="rent_editor")
                if edited_df is not None: st.session_state['rent_roll_data'] = edited_df
                
                df_calc = pd.DataFrame(st.session_state['rent_roll_data'])
                sum_dep = pd.to_numeric(df_calc['보증금'], errors='coerce').fillna(0).sum()
                sum_rent = pd.to_numeric(df_calc['임대료'], errors='coerce').fillna(0).sum()
                sum_maint = pd.to_numeric(df_calc['관리비'], errors='coerce').fillna(0).sum()
                
                st.subheader("💰 금액 정보")
                r1, r2, r3 = st.columns(3)
                deposit_val = comma_input("보증금", "만원", "deposit", int(sum_dep))
                rent_val = comma_input("월임대료", "만원", "rent", int(sum_rent))
                maint_val = comma_input("관리비", "만원", "maint", int(sum_maint))
                
                r4, r5, r6 = st.columns(3)
                loan_val = comma_input("융자금", "억원", "loan", 0)
                price_val = comma_input("매매금액", "억원", "price", 0)
                
                try: yield_rate = ((rent_val * 12) / ((price_val * 10000) - deposit_val)) * 100 if ((price_val * 10000) - deposit_val) > 0 else 0
                except: yield_rate = 0
                r6.metric("수익률", f"{yield_rate:.2f}%")
                
                finance_data = {"price": price_val, "deposit": deposit_val, "rent": rent_val, "maintenance": maint_val, "loan": loan_val, "yield": yield_rate, "land_pyeong_price_val": (price_val*100000000)/(info['platArea']*0.3025)/10000 if info['platArea'] else 0}
                
                st.subheader("🔍 AI 물건분석")
                env_options = ["역세권", "광대로변", "먹자상권", "랜드마크", "급매물", "주차편리", "감정가이하", "초역세권", "대로변", "오피스상권", "법조타운", "사옥추천", "수려한외관", "용적률이득", "더블역세권", "대로코너", "항아리상권", "핫플레이스", "수익형", "신축빌딩", "신축부지용", "트리플역세권", "이면코너", "학군지", "메디컬입지", "시세차익", "관리상태최상", "명도완료", "광역환승", "이면초입", "숲세권", "가시성우수", "밸류업유망", "리모델링", "명도협의가능"]
                cols = st.columns(7)
                selected_envs = [opt for i, opt in enumerate(env_options) if cols[i%7].checkbox(opt)]
                
                comp_file = st.file_uploader("비교 분석 엑셀", type=['xlsx'])
                comp_df = pd.read_excel(comp_file) if comp_file else None
                user_comment = st.text_area("추가 특징")
                
                if st.button("🤖 인사이트요약"):
                    st.session_state['generated_candidates'] = generate_insight_candidates(info, finance_data, st.session_state['zoning'], selected_envs, user_comment, comp_df)
                    st.rerun()
                
                if st.session_state['generated_candidates']:
                    for cand in st.session_state['generated_candidates']:
                        if st.checkbox(cand, key=cand):
                            if cand not in st.session_state['final_selected_insights']: st.session_state['final_selected_insights'].append(cand)
                            st.session_state['generated_candidates'].remove(cand); st.rerun()
                            
                st.write("최종 선택된 포인트:", st.session_state['final_selected_insights'])
                
                st.subheader("📥 저장")
                c_p, c_x = st.columns(2)
                ppt_9 = c_p.file_uploader("9장 템플릿", type=['pptx'], key="t9")
                ppt_1 = c_p.file_uploader("1장 템플릿", type=['pptx'], key="t1")
                
                if ppt_9 or ppt_1:
                    ppt_data = create_pptx(info, location['full_addr'], finance_data, st.session_state['zoning'], location['lat'], location['lng'], land_price, st.session_state['final_selected_insights'], images_map, pd.DataFrame(st.session_state['rent_roll_data']), ppt_9, ppt_1)
                    c_p.download_button("PPT 다운로드", ppt_data, "report.pptx")
                    
                xls_data = create_excel(info, location['full_addr'], finance_data, st.session_state['zoning'], location['lat'], location['lng'], land_price, st.session_state['final_selected_insights'], images_map.get('u2'))
                c_x.download_button("엑셀 다운로드", xls_data, "report.xlsx")
